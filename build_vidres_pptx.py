#!/usr/bin/env python3
"""Gera a apresentação estratégica Vidres Portugal (.pptx) em português."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- Paleta de marca ----
NAVY = RGBColor(0x0E, 0x2A, 0x47)      # azul-escuro
TEAL = RGBColor(0x12, 0x8C, 0x7E)      # verde-azulado (acento)
ORANGE = RGBColor(0xE3, 0x6A, 0x1E)    # laranja (destaque)
LIGHT = RGBColor(0xF2, 0xF5, 0xF7)     # fundo claro
GREY = RGBColor(0x55, 0x66, 0x70)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x24, 0x2B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, l, t, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.0):
    """runs: lista de parágrafos; cada parágrafo é lista de (texto, size, bold, color)."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for (s, size, bold, color) in para:
            r = p.add_run()
            r.text = s
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = "Calibri"
    return tb


def accent_bar(slide, color=ORANGE):
    box(slide, 0, 0, Inches(0.18), SH, color)


def footer(slide, n, dark_bg=False):
    c = WHITE if dark_bg else GREY
    text(slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.4),
         [[("Vidres Portugal — Manual Estratégico", 9, False, c)]])
    text(slide, Inches(11.3), Inches(7.05), Inches(1.8), Inches(0.4),
         [[(f"{n:02d}", 9, True, c)]], align=PP_ALIGN.RIGHT)


def section_header(slide, kicker, title, dark_bg=False):
    accent_bar(slide)
    kc = TEAL if not dark_bg else ORANGE
    tc = NAVY if not dark_bg else WHITE
    text(slide, Inches(0.6), Inches(0.45), Inches(12), Inches(0.4),
         [[(kicker.upper(), 13, True, kc)]])
    text(slide, Inches(0.6), Inches(0.85), Inches(12.1), Inches(0.9),
         [[(title, 30, True, tc)]])


def bullets(slide, l, t, w, h, items, size=15, color=DARK, gap=8):
    runs = []
    for it in items:
        runs.append([("•  ", size, True, TEAL), (it, size, False, color)])
    text(slide, l, t, w, h, runs, space_after=gap, line_spacing=1.05)


# =====================================================================
# SLIDE 1 — Capa
# =====================================================================
s = add_slide()
bg(s, NAVY)
box(s, 0, Inches(5.55), SW, Inches(0.12), ORANGE)
box(s, 0, Inches(5.67), SW, Inches(0.06), TEAL)
text(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.5),
     [[("VIDRES PORTUGAL", 20, True, TEAL)]])
text(s, Inches(0.9), Inches(2.05), Inches(11.6), Inches(2.2),
     [[("Manual Estratégico", 52, True, WHITE)],
      [("O Parceiro Ágil de Inovação Técnica", 30, True, ORANGE)],
      [("da Cerâmica Portuguesa", 30, True, ORANGE)]],
     space_after=4)
text(s, Inches(0.9), Inches(5.9), Inches(11.5), Inches(1.2),
     [[("Matéria-prima para loiça & cerâmica  ·  Estratégia competitiva contra Torrecid / Esmalglass / Ferro",
        14, False, RGBColor(0xC8, 0xD4, 0xDC))],
      [("Vista Alegre  ·  Costa Verde  ·  Cluster de Aveiro", 13, True, WHITE)]],
     space_after=4)

# =====================================================================
# SLIDE 2 — Tese central
# =====================================================================
s = add_slide()
bg(s, LIGHT)
accent_bar(s)
text(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.4),
     [[("TESE CENTRAL", 13, True, TEAL)]])
box(s, Inches(0.9), Inches(1.4), Inches(11.5), Inches(2.1), NAVY)
text(s, Inches(1.3), Inches(1.65), Inches(10.7), Inches(1.7),
     [[("A Vidres Portugal deve tornar-se o parceiro ágil de inovação técnica dos fabricantes de cerâmica portugueses — e não um fornecedor de matéria-prima em regime de commodity.",
        23, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(0.9), Inches(3.95), Inches(11.5), Inches(0.6),
     [[("Não vencemos a Torrecid em escala. Vencemo-la por sermos:", 17, True, NAVY)]])

labels = [("MAIS RÁPIDOS", "amostra à medida\n≤ 10 dias úteis"),
          ("MAIS PRÓXIMOS", "engenheiro na linha\nem 24–48h"),
          ("MAIS FLEXÍVEIS", "lotes-piloto\n50–200 kg"),
          ("MAIS SUSTENTÁVEIS", "8–12% poupança\nde energia")]
x = Inches(0.9)
cw = Inches(2.75)
gap = Inches(0.13)
for i, (h, d) in enumerate(labels):
    bx = Emu(int(x) + i * (int(cw) + int(gap)))
    box(s, bx, Inches(4.75), cw, Inches(1.7), TEAL if i % 2 == 0 else ORANGE)
    text(s, bx, Inches(4.95), cw, Inches(0.6),
         [[(h, 15, True, WHITE)]], align=PP_ALIGN.CENTER)
    text(s, bx, Inches(5.55), cw, Inches(0.85),
         [[(line, 12, False, WHITE)] for line in d.split("\n")],
         align=PP_ALIGN.CENTER, space_after=2)
footer(s, 2)

# =====================================================================
# SLIDE 3 — Mercado
# =====================================================================
s = add_slide()
bg(s, WHITE)
section_header(s, "1. Mercado", "Mercado português de matérias-primas cerâmicas")
# três segmentos
segs = [("LOIÇA PREMIUM", "Vista Alegre · Costa Verde · Spal",
         "Qualidade > preço.\nA consistência do lote domina.", ORANGE),
        ("VOLUME / GRÉS", "Matceramica · Grestel · regionais",
         "Sensível a preço + energia,\nmas rejeita defeitos.", TEAL),
        ("SANITÁRIO / AZULEJO", "Sanindusa · Revigrés · Love Tiles",
         "Guiado por volume\ne por energia.", NAVY)]
x = Inches(0.6)
cw = Inches(3.95)
gap = Inches(0.2)
for i, (h, ex, d, col) in enumerate(segs):
    bx = Emu(int(x) + i * (int(cw) + int(gap)))
    box(s, bx, Inches(1.95), cw, Inches(2.3), col)
    text(s, bx, Inches(2.15), cw, Inches(0.5), [[(h, 15, True, WHITE)]], align=PP_ALIGN.CENTER)
    text(s, bx, Inches(2.7), cw, Inches(0.5), [[(ex, 11, False, WHITE)]], align=PP_ALIGN.CENTER)
    text(s, bx, Inches(3.25), cw, Inches(0.9),
         [[(line, 12, True, WHITE)] for line in d.split("\n")], align=PP_ALIGN.CENTER, space_after=2)
box(s, Inches(0.6), Inches(4.55), Inches(12.1), Inches(2.05), LIGHT)
text(s, Inches(0.9), Inches(4.7), Inches(11.5), Inches(0.5),
     [[("Cluster de Aveiro / Centro — a vantagem estrutural", 17, True, NAVY)]])
bullets(s, Inches(0.9), Inches(5.25), Inches(11.5), Inches(1.3), [
    "Aveiro / Ílhavo / Vagos — porcelana e loiça (Vista Alegre, Costa Verde, Grestel).",
    "Polo técnico local = 60–90 min das contas-chave. Torrecid/Esmalglass: ~900 km de Castellón.",
    "A proximidade é o fosso defensivo mais forte que temos.",
], size=14, gap=6)
footer(s, 3)

# =====================================================================
# SLIDE 4 — Panorama competitivo
# =====================================================================
s = add_slide()
bg(s, WHITE)
section_header(s, "2. Panorama competitivo", "A lacuna estratégica: ninguém ocupa o quadrante vencedor")
comp = [("Torrecid", "Escala global, laboratórios fortes",
         "Lento em pequenos lotes; longe de Portugal; rígido"),
        ("Esmalglass-Itaca", "I&D profunda, tintas digitais",
         "Decisão corporativa lenta; Portugal à distância"),
        ("Ferro / Vibrantz", "Ciência da cor, sem chumbo",
         "Reorganização pós-fusão; pouca intimidade ibérica"),
        ("Locais / pequenos", "Baratos, locais, flexíveis",
         "I&D fraca, sem roteiro, lotes inconsistentes")]
y = Inches(1.85)
for i, (n, f, w) in enumerate(comp):
    ty = Emu(int(y) + i * int(Inches(0.82)))
    box(s, Inches(0.6), ty, Inches(2.5), Inches(0.7), NAVY)
    text(s, Inches(0.6), ty, Inches(2.5), Inches(0.7), [[(n, 14, True, WHITE)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    box(s, Inches(3.2), ty, Inches(3.6), Inches(0.7), RGBColor(0xDD, 0xEC, 0xEA))
    text(s, Inches(3.35), ty, Inches(3.4), Inches(0.7), [[(f, 12, False, DARK)]],
         anchor=MSO_ANCHOR.MIDDLE)
    box(s, Inches(6.9), ty, Inches(5.8), Inches(0.7), RGBColor(0xFB, 0xE6, 0xD8))
    text(s, Inches(7.05), ty, Inches(5.6), Inches(0.7),
         [[("Fraqueza: ", 12, True, ORANGE), (w, 12, False, DARK)]], anchor=MSO_ANCHOR.MIDDLE)
box(s, Inches(0.6), Inches(5.35), Inches(12.1), Inches(1.25), TEAL)
text(s, Inches(0.9), Inches(5.5), Inches(11.5), Inches(1.0),
     [[("◆  A posição vencedora da Vidres = alta tecnologia + alta proximidade.", 18, True, WHITE)],
      [("Os espanhóis têm tecnologia mas não proximidade; os locais têm proximidade mas não tecnologia.",
        14, False, WHITE)]], space_after=4)
footer(s, 4)

# =====================================================================
# SLIDE 5 — Pontos de dor
# =====================================================================
s = add_slide()
bg(s, LIGHT)
section_header(s, "3. Pontos de dor", "As 4 coisas que as fábricas realmente compram")
pains = [("1 · Consistência entre lotes",
          "O principal motor de abandono. Tom, fusão, brilho ou defeitos que variam = cliente perdido.",
          "SPC + Certificado de Análise · ΔE ≤ 0,5 · SLA contratual"),
         ("2 · Eficiência energética",
          "Gás e eletricidade são a rubrica #1 do P&L. Querem menor temperatura e ciclo mais curto.",
          "Eco-Low Fire: −30 a −50 °C → 8–12% de energia"),
         ("3 · Desenvolvimento à medida rápido",
          "O laboratório da Torrecid é a sua fortaleza. As fábricas precisam de respostas rápidas.",
          "Laboratório local + cor digital → amostra ≤ 10 dias"),
         ("4 · Apoio técnico no local",
          "A proximidade vence Espanha. Um engenheiro que resolve a linha em 24–48h é decisivo.",
          "Engenheiros de campo em Portugal · resposta 24–48h")]
x = Inches(0.6)
cw = Inches(5.95)
ch = Inches(2.15)
gx = Inches(0.2)
gy = Inches(0.25)
for i, (h, d, ans) in enumerate(pains):
    col = i % 2
    row = i // 2
    bx = Emu(int(x) + col * (int(cw) + int(gx)))
    by = Emu(int(Inches(1.95)) + row * (int(ch) + int(gy)))
    box(s, bx, by, cw, ch, WHITE)
    box(s, bx, by, Inches(0.12), ch, ORANGE)
    text(s, Emu(int(bx) + int(Inches(0.3))), Emu(int(by) + int(Inches(0.15))),
         Emu(int(cw) - int(Inches(0.5))), Inches(0.5), [[(h, 16, True, NAVY)]])
    text(s, Emu(int(bx) + int(Inches(0.3))), Emu(int(by) + int(Inches(0.65))),
         Emu(int(cw) - int(Inches(0.5))), Inches(0.85), [[(d, 12.5, False, DARK)]], line_spacing=1.05)
    text(s, Emu(int(bx) + int(Inches(0.3))), Emu(int(by) + int(Inches(1.55))),
         Emu(int(cw) - int(Inches(0.5))), Inches(0.5),
         [[("→ ", 12.5, True, TEAL), (ans, 12.5, True, TEAL)]])
footer(s, 5)

# =====================================================================
# SLIDE 6 — Otimização de custos
# =====================================================================
s = add_slide()
bg(s, WHITE)
section_header(s, "4. Custos", "Três alavancas para proteger a margem")
levers = [("A · Abastecimento ibérico local",
           ["Sílica, feldspato e casco reciclado de origem ibérica",
            "Recuperar casco dos rejeitos do próprio cliente",
            "↓ transporte · ↓ prazo · ↓ CO₂ · economia circular"]),
          ("B · 5 fritas-base padronizadas",
           ["Não 50 receitas — 5 bases + aditivos/corantes",
            "Poder de compra e QC mais simples",
            "Resposta à medida mais rápida e consistente"]),
          ("C · Correspondência de cor digital",
           ["Formulação automatizada por espectrofotómetro",
            "↓ desperdício · ↓ ciclos de teste",
            "↓ tempo e custo de desenvolvimento"])]
x = Inches(0.6)
cw = Inches(3.95)
gap = Inches(0.2)
cols = [TEAL, NAVY, ORANGE]
for i, (h, items) in enumerate(levers):
    bx = Emu(int(x) + i * (int(cw) + int(gap)))
    box(s, bx, Inches(1.95), cw, Inches(0.75), cols[i])
    text(s, bx, Inches(1.95), cw, Inches(0.75), [[(h, 15, True, WHITE)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    box(s, bx, Inches(2.7), cw, Inches(2.5), LIGHT)
    bullets(s, Emu(int(bx) + int(Inches(0.25))), Inches(2.95),
            Emu(int(cw) - int(Inches(0.45))), Inches(2.2), items, size=12.5, gap=8)
box(s, Inches(0.6), Inches(5.5), Inches(12.1), Inches(1.1), NAVY)
text(s, Inches(0.9), Inches(5.65), Inches(11.5), Inches(0.85),
     [[("Efeito combinado A+B+C: ", 17, True, ORANGE),
       ("redução de 20–35% do custo-de-servir e de I&D —", 17, True, WHITE)],
      [("financia preço competitivo E serviço técnico gratuito que fideliza as contas.", 14, False, WHITE)]],
     space_after=3)
footer(s, 6)

# =====================================================================
# SLIDE 7 — Roteiro de inovação
# =====================================================================
s = add_slide()
bg(s, NAVY)
accent_bar(s, ORANGE)
text(s, Inches(0.6), Inches(0.45), Inches(12), Inches(0.4),
     [[("5. ROTEIRO DE INOVAÇÃO", 13, True, ORANGE)]])
text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.8),
     [[("Duas bandeiras: de fornecedor a parceiro tecnológico", 28, True, WHITE)]])
# Bandeira 1
box(s, Inches(0.6), Inches(1.9), Inches(5.95), Inches(3.05), TEAL)
text(s, Inches(0.85), Inches(2.05), Inches(5.5), Inches(0.5),
     [[("BANDEIRA #1 — arma comercial", 12, True, WHITE)]])
text(s, Inches(0.85), Inches(2.5), Inches(5.5), Inches(0.7),
     [[("Sistema Eco-Low Fire Smart Frit", 20, True, WHITE)]])
bullets(s, Inches(0.85), Inches(3.3), Inches(5.5), Inches(1.6), [
    "Baixa a temperatura de cozedura 30–50 °C",
    "Ciclo mais curto, menos consumo de energia",
    "Mantém brancura e brilho (loiça premium)",
], size=13, color=WHITE, gap=8)
box(s, Inches(0.85), Inches(4.35), Inches(5.45), Inches(0.45), WHITE)
text(s, Inches(0.85), Inches(4.35), Inches(5.45), Inches(0.45),
     [[("8–12% de poupança de energia — mais forte que o preço", 13, True, TEAL)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# Bandeira 2
box(s, Inches(6.75), Inches(1.9), Inches(5.95), Inches(3.05), ORANGE)
text(s, Inches(7.0), Inches(2.05), Inches(5.5), Inches(0.5),
     [[("BANDEIRA #2 — posicionamento tech", 12, True, WHITE)]])
text(s, Inches(7.0), Inches(2.5), Inches(5.5), Inches(0.7),
     [[("Plataforma de Previsão de Vidrados por IA", 20, True, WHITE)]])
bullets(s, Inches(7.0), Inches(3.35), Inches(5.5), Inches(1.6), [
    "Input: pasta + curva de forno + acabamento",
    "Recomenda a receita de partida e prevê defeitos",
    "Fosso de dados → custo de mudança sobe",
], size=13, color=WHITE, gap=8)
# faixa de sequenciamento
box(s, Inches(0.6), Inches(5.2), Inches(12.1), Inches(1.4), RGBColor(0x13, 0x35, 0x55))
text(s, Inches(0.85), Inches(5.32), Inches(11.6), Inches(0.4),
     [[("Sequenciamento — primeiro a vitória energética, depois a IA", 14, True, ORANGE)]])
seq = "H1 (M1–6) 5 fritas-base + ΔE   ▶   H2 (M4–12) Eco-Low-Fire v1 + pilotos   ▶   H3 (M9–18) cor digital + v2   ▶   H4 (M12–24) IA → contas-chave"
text(s, Inches(0.85), Inches(5.78), Inches(11.6), Inches(0.7),
     [[(seq, 12.5, False, WHITE)]])
footer(s, 7, dark_bg=True)

# =====================================================================
# SLIDE 8 — Portfólio de 5 inovações de matéria-prima (grelha)
# =====================================================================
s = add_slide()
bg(s, WHITE)
section_header(s, "5A. Portfólio de inovação",
               "5 inovações de matéria-prima — não \"o mesmo, mais barato\"")
innov = [
    ("1", "Porcelana\nTranslúcida Vegan", "Fine china SEM osso",
     "Fosfato de cálcio sintético + wollastonite",
     "Premium · vegan · ESG", ORANGE),
    ("2", "BioShield", "Vidrado antimicrobiano",
     "Iões Ag / Zn / Cu · sem chumbo/cádmio",
     "Auto-higienizante (hospitality/saúde)", TEAL),
    ("3", "TerraCircular", "Matéria-prima circular PT",
     "Cinza de cortiça · biossílica de arroz · conchas · cullet",
     "−CO₂ + cobertura CBAM", NAVY),
    ("4", "ChromaSafe", "Cor vibrante segura",
     "Pigmentos encapsulados em zircão + terras-raras",
     "Vermelhos food-safe a baixa cozedura", RGBColor(0x6A, 0x4C, 0x93)),
    ("5", "DuraGlaze /\nSensaGlaze", "Superfícies funcionais",
     "CTE ajustado + nano (zircónia) + termocrómico",
     "Durabilidade + efeito \"wow\"", RGBColor(0xB0, 0x3A, 0x2E)),
]
x0 = Inches(0.6)
cw = Inches(2.25)
gap = Inches(0.135)
top = Inches(1.95)
ch = Inches(4.05)
for i, (num, name, sub, base, benefit, col) in enumerate(innov):
    bx = Emu(int(x0) + i * (int(cw) + int(gap)))
    box(s, bx, top, cw, ch, LIGHT)
    box(s, bx, top, cw, Inches(1.15), col)
    text(s, bx, Emu(int(top) + int(Inches(0.08))), cw, Inches(0.4),
         [[(num, 22, True, WHITE)]], align=PP_ALIGN.CENTER)
    text(s, bx, Emu(int(top) + int(Inches(0.48))), cw, Inches(0.7),
         [[(line, 13, True, WHITE)] for line in name.split("\n")],
         align=PP_ALIGN.CENTER, space_after=0, line_spacing=0.95)
    pad = int(Inches(0.18))
    iw = Emu(int(cw) - 2 * pad)
    text(s, Emu(int(bx) + pad), Emu(int(top) + int(Inches(1.3))), iw, Inches(0.5),
         [[(sub, 11.5, True, col)]], align=PP_ALIGN.CENTER, line_spacing=0.95)
    text(s, Emu(int(bx) + pad), Emu(int(top) + int(Inches(1.95))), iw, Inches(1.3),
         [[("Componentes:", 10, True, GREY)], [(base, 10.5, False, DARK)]],
         align=PP_ALIGN.CENTER, space_after=2, line_spacing=1.0)
    text(s, Emu(int(bx) + pad), Emu(int(top) + int(Inches(3.25))), iw, Inches(0.7),
         [[(benefit, 11, True, NAVY)]], align=PP_ALIGN.CENTER, line_spacing=0.95)
box(s, Inches(0.6), Inches(6.15), Inches(12.1), Inches(0.62), DARK)
text(s, Inches(0.6), Inches(6.15), Inches(12.1), Inches(0.62),
     [[("Vendemos componentes que ninguém mais tem — segurança alimentar, circularidade, durabilidade e design funcional.",
        13, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 8)

# =====================================================================
# SLIDE 9 — Priorização do portfólio de inovação
# =====================================================================
s = add_slide()
bg(s, LIGHT)
section_header(s, "5A. Portfólio de inovação", "Priorização — o que lançar primeiro e porquê")
rows = [
    ("Inovação", "Dor que resolve", "Margem", "I&D", "Prioridade", True),
    ("3 · TerraCircular", "Custo + CO₂ + CBAM", "Média", "Baixo", "IMEDIATA", False),
    ("2 · BioShield", "Higiene / valor funcional", "Alta", "Médio", "Alta", False),
    ("1 · Porcelana Vegan", "Diferenciação premium / ESG", "Alta", "Médio", "Alta (Vista Alegre)", False),
    ("4 · ChromaSafe", "Cor segura difícil", "Alta", "Médio-alto", "Média", False),
    ("5 · DuraGlaze/SensaGlaze", "Durabilidade + design", "Média-alta", "Alto", "Faseada", False),
]
colx = [Inches(0.6), Inches(3.7), Inches(7.4), Inches(8.7), Inches(10.0)]
colw = [Inches(3.0), Inches(3.6), Inches(1.2), Inches(1.2), Inches(2.7)]
ry = Inches(1.95)
rh = Inches(0.62)
for r, row in enumerate(rows):
    ty = Emu(int(ry) + r * int(rh))
    is_head = row[5]
    rowcol = NAVY if is_head else (WHITE if r % 2 else RGBColor(0xE7, 0xEE, 0xF1))
    box(s, Inches(0.6), ty, Inches(12.1), rh, rowcol)
    for c in range(5):
        tcol = WHITE if is_head else (ORANGE if (c == 4 and r == 1) else DARK)
        bold = is_head or (c == 0) or (c == 4 and r == 1)
        text(s, Emu(int(colx[c]) + int(Inches(0.1))), ty, colw[c], rh,
             [[(row[c], 12.5 if is_head else 12, bold, tcol)]], anchor=MSO_ANCHOR.MIDDLE)
box(s, Inches(0.6), Inches(5.75), Inches(12.1), Inches(1.0), TEAL)
text(s, Inches(0.9), Inches(5.9), Inches(11.5), Inches(0.75),
     [[("Sequência: ", 15, True, WHITE),
       ("TerraCircular + BioShield já  →  Porcelana Vegan como gancho na Vista Alegre  →  ChromaSafe e DuraGlaze sobre a base Eco-Low-Fire.",
        15, False, WHITE)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
footer(s, 9)

# =====================================================================
# SLIDE 10 — Go-to-market
# =====================================================================
s = add_slide()
bg(s, WHITE)
section_header(s, "6. Entrada nas contas-chave", "Jogo de contas nomeadas — não de volume")
tiers = [("NÍVEL 1 — Farol", "Vista Alegre · Costa Verde", "Piloto Eco-Low-Fire numa linha", ORANGE),
         ("NÍVEL 2 — Volume", "Grestel · Matceramica · Spal", "SLA de consistência + energia", TEAL),
         ("NÍVEL 3 — Expansão", "Sanindusa · Revigrés · Love Tiles", "Engobes/vidrados + serviço local", NAVY)]
for i, (h, ex, w, col) in enumerate(tiers):
    ty = Emu(int(Inches(1.95)) + i * int(Inches(0.78)))
    box(s, Inches(0.6), ty, Inches(2.7), Inches(0.65), col)
    text(s, Inches(0.6), ty, Inches(2.7), Inches(0.65), [[(h, 13, True, WHITE)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(3.5), ty, Inches(4.0), Inches(0.65),
         [[(ex, 13, True, NAVY)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(7.6), ty, Inches(5.1), Inches(0.65),
         [[(w, 12.5, False, DARK)]], anchor=MSO_ANCHOR.MIDDLE)
box(s, Inches(0.6), Inches(4.5), Inches(12.1), Inches(2.1), LIGHT)
text(s, Inches(0.9), Inches(4.65), Inches(11.5), Inches(0.5),
     [[("A jogada de \"aterragem\" em 5 passos", 17, True, NAVY)]])
steps = "1. Visita de diagnóstico gratuita (relatório de oportunidade)   →   2. Piloto numa única linha (baixo risco)   →   3. Resultado quantificado em €/ano   →   4. Contrato com SLA de consistência   →   5. Expandir + plataforma de IA = fidelização"
text(s, Inches(0.9), Inches(5.2), Inches(11.5), Inches(1.3),
     [[(steps, 14, False, DARK)]], line_spacing=1.2)
footer(s, 10)

# =====================================================================
# SLIDE 11 — Pilotos
# =====================================================================
s = add_slide()
bg(s, LIGHT)
section_header(s, "7. Projetos-piloto", "Dois pilotos farol = prova quantificada e inegável")
# Piloto A
box(s, Inches(0.6), Inches(1.95), Inches(5.95), Inches(4.4), WHITE)
box(s, Inches(0.6), Inches(1.95), Inches(5.95), Inches(0.8), NAVY)
text(s, Inches(0.85), Inches(1.95), Inches(5.5), Inches(0.8),
     [[("PILOTO A — Vista Alegre", 18, True, WHITE)],
      [("Porcelana premium · a barra de qualidade mais alta", 11, False, RGBColor(0xC8,0xD4,0xDC))]],
     anchor=MSO_ANCHOR.MIDDLE, space_after=2)
bullets(s, Inches(0.85), Inches(3.0), Inches(5.5), Inches(3.2), [
    "Provar Eco-Low-Fire sem perda de brancura/brilho",
    "−30 a −50 °C na temperatura de pico",
    "≥ 8% de redução medida de energia",
    "ΔE ≤ 0,5 em 3 lotes consecutivos",
    "Vencer aqui = certificado de qualidade para todas as contas",
], size=13.5, gap=10)
# Piloto B
box(s, Inches(6.75), Inches(1.95), Inches(5.95), Inches(4.4), WHITE)
box(s, Inches(6.75), Inches(1.95), Inches(5.95), Inches(0.8), TEAL)
text(s, Inches(7.0), Inches(1.95), Inches(5.5), Inches(0.8),
     [[("PILOTO B — Costa Verde", 18, True, WHITE)],
      [("Louça de volume · prova de serviço e SLA", 11, False, WHITE)]],
     anchor=MSO_ANCHOR.MIDDLE, space_after=2)
bullets(s, Inches(7.0), Inches(3.0), Inches(5.5), Inches(3.2), [
    "SLA de consistência + serviço no local 24–48h",
    "Taxa de defeitos ↓ (alvo −20%)",
    "Energia/peça ↓ 8–12%",
    "Resposta do engenheiro ≤ 48h registada",
    "Poupança em €/ano validada pela equipa de processo",
], size=13.5, gap=10)
footer(s, 11)

# =====================================================================
# SLIDE 12 — Plano a 24 meses
# =====================================================================
s = add_slide()
bg(s, WHITE)
section_header(s, "8. Execução", "Plano de execução a 24 meses")
phases = [("FASE 1\nM0–6\nFUNDAÇÃO",
           ["Laboratório + armazém em Aveiro", "2–3 engenheiros de campo",
            "5 fritas-base + SPC/ΔE", "Assinar pilotos VA + CV"], TEAL),
          ("FASE 2\nM4–12\nPROVA",
           ["Eco-Low-Fire v1 validado", "Executar Pilotos A e B",
            "Cor digital ativa", "1ºs casos de estudo"], NAVY),
          ("FASE 3\nM9–18\nESCALA",
           ["Pilotos → contratos SLA", "Eco-Low-Fire v2",
            "Expandir Nível 2", "Abordar Nível 3"], ORANGE),
          ("FASE 4\nM12–24\nFIDELIZAÇÃO",
           ["Plataforma de IA (beta→rollout)", "Fosso de dados",
            "Co-desenvolvimento", "Parceiro por defeito do cluster"], RGBColor(0x6A,0x4C,0x93))]
x = Inches(0.6)
cw = Inches(2.95)
gap = Inches(0.13)
for i, (h, items, col) in enumerate(phases):
    bx = Emu(int(x) + i * (int(cw) + int(gap)))
    box(s, bx, Inches(1.95), cw, Inches(1.15), col)
    text(s, bx, Inches(1.95), cw, Inches(1.15),
         [[(line, 14 if j == 0 else (11 if j == 1 else 14), True, WHITE)]
          for j, line in enumerate(h.split("\n"))],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=1)
    box(s, bx, Inches(3.1), cw, Inches(3.0), LIGHT)
    bullets(s, Emu(int(bx) + int(Inches(0.2))), Inches(3.3),
            Emu(int(cw) - int(Inches(0.35))), Inches(2.7), items, size=12, gap=10)
footer(s, 12)

# =====================================================================
# SLIDE 13 — KPIs
# =====================================================================
s = add_slide()
bg(s, LIGHT)
section_header(s, "Anexo A", "Painel de KPIs — acompanhar desde o dia 1")
kpis = [("ΔE entre lotes", "≤ 0,5"), ("Redução de energia", "8–12%"),
        ("Prazo de amostra", "≤ 10 dias"), ("Resposta no local", "24–48h"),
        ("Conversão piloto→contrato", "≥ 60%"), ("Custo-de-servir / I&D", "−20–35%")]
x = Inches(0.6)
cw = Inches(3.85)
ch = Inches(1.85)
gx = Inches(0.2)
gy = Inches(0.25)
cols = [TEAL, ORANGE, NAVY]
for i, (k, v) in enumerate(kpis):
    col = i % 3
    row = i // 3
    bx = Emu(int(x) + col * (int(cw) + int(gx)))
    by = Emu(int(Inches(2.1)) + row * (int(ch) + int(gy)))
    box(s, bx, by, cw, ch, WHITE)
    box(s, bx, by, cw, Inches(0.1), cols[col])
    text(s, bx, Emu(int(by) + int(Inches(0.3))), cw, Inches(0.9),
         [[(v, 34, True, cols[col])]], align=PP_ALIGN.CENTER)
    text(s, bx, Emu(int(by) + int(Inches(1.25))), cw, Inches(0.5),
         [[(k, 14, True, NAVY)]], align=PP_ALIGN.CENTER)
footer(s, 13)

# =====================================================================
# SLIDE 14 — Resumo / pitch de elevador
# =====================================================================
s = add_slide()
bg(s, NAVY)
box(s, 0, Inches(2.9), SW, Inches(0.1), ORANGE)
text(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.5),
     [[("RESUMO — PITCH DE ELEVADOR", 14, True, TEAL)]])
text(s, Inches(0.9), Inches(1.25), Inches(11.5), Inches(1.6),
     [[("Tornar-nos o parceiro ágil de inovação técnica da cerâmica portuguesa.",
        30, True, WHITE)]])
text(s, Inches(0.9), Inches(3.25), Inches(11.5), Inches(3.3),
     [[("Vencemos não por superar a Torrecid, a Esmalglass e a Ferro em escala, mas por ocupar o quadrante vazio: alta tecnologia + alta proximidade local.",
        16, False, RGBColor(0xD8,0xE2,0xE8))],
      [("Entregamos as 4 coisas que as fábricas compram — consistência de lote, faturas de energia mais baixas, desenvolvimento à medida rápido e um engenheiro na linha em 24–48h.",
        16, False, RGBColor(0xD8,0xE2,0xE8))],
      [("O Eco-Low Fire Smart Frit transforma a maior dor (energia) no nosso argumento mais forte (8–12%). A Plataforma de IA torna a relação num fosso de dados.",
        16, False, RGBColor(0xD8,0xE2,0xE8))],
      [("Provamo-lo na Vista Alegre e na Costa Verde — e escalamos pelo cluster em 24 meses.",
        17, True, ORANGE)]], space_after=12, line_spacing=1.05)
footer(s, 14, dark_bg=True)

out = "Vidres_Portugal_Apresentacao_Estrategica.pptx"
prs.save(out)
print(f"OK -> {out} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
